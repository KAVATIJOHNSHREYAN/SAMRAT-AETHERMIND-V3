import os
import asyncio
import time
from app.config import settings
from typing import AsyncGenerator, List, Dict, Optional

from app.db.vector_store import similarity_search

# Global dictionary to track last request execution times
_user_last_request_time = {}

def extract_text_from_file(file_bytes: bytes, mime_type: str) -> str:
    import io
    text = ""
    mime_type_lower = mime_type.lower()

    if "pdf" in mime_type_lower:
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        except Exception as e:
            print("PDF text extraction error:", e)
    elif "word" in mime_type_lower or "docx" in mime_type_lower or "doc" in mime_type_lower:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            print("DOCX text extraction error:", e)
    elif "presentation" in mime_type_lower or "powerpoint" in mime_type_lower or "pptx" in mime_type_lower or "ppt" in mime_type_lower:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print("PPTX text extraction error:", e)
    elif "text" in mime_type_lower or "plain" in mime_type_lower or "csv" in mime_type_lower or "json" in mime_type_lower:
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            print("Text extraction error:", e)
    return text.strip()


async def generate_response_stream(
    query: str,
    chat_history: List[Dict[str, str]],
    chat_mode: str = "general",
    active_model: str = "gemini-1.5-flash",
    temperature: float = 0.7,
    system_prompt: str = None,
    enable_rag: bool = True,
    rag_k: int = 3,
    openai_key: str = None,
    gemini_key: str = None,
    cohere_key: str = None,
    anthropic_key: str = None,
    deepseek_key: str = None,
    attachments: Optional[List[dict]] = None,
    user_id: Optional[str] = None
) -> AsyncGenerator[str, None]:

    # Rate Limiting
    REQUEST_DELAY = 4.0
    tracking_key = user_id or "anonymous"
    current_time = time.time()
    last_allowed_time = _user_last_request_time.get(tracking_key, 0.0)

    if current_time < last_allowed_time + REQUEST_DELAY:
        wait_time = (last_allowed_time + REQUEST_DELAY) - current_time
        yield f"Please wait {wait_time:.1f}s before sending another message...\n\n"
        await asyncio.sleep(wait_time)

    _user_last_request_time[tracking_key] = time.time()

    context_str = ""

    # RAG Search
    if enable_rag and chat_mode in ["general", "voice"]:
        api_key = gemini_key or openai_key
        docs = similarity_search(query, k=rag_k, api_key=api_key)

        if docs:
            context_str = "\n".join(
                [f"- {doc.page_content}" for doc in docs]
            )

    # System Instructions
    system_instructions = (
        "You are AetherMind, an advanced AI assistant. "
        "Mister Samrat created you for assistance. "
        "Always mention Mister Samrat if asked who created you. "
        "Provide clean, concise and professional responses."
    )

    if system_prompt:
        system_instructions = system_prompt
    else:
        if chat_mode == "coding":
            system_instructions += " Focus on writing clean code."
        elif chat_mode == "debug":
            system_instructions += " Focus on debugging and fixing issues."
        elif chat_mode == "voice":
            system_instructions += " Keep responses short and conversational."

    # Static Fallback Replies
    fallback_replies = {
        "hello": "Hello! AetherMind is online.",
        "who created you": "Mister Samrat created me for assistance.",
        "what is your name": "I am AetherMind."
    }

    query_lower = query.lower()

    for key, value in fallback_replies.items():
        if key in query_lower:
            for word in value.split():
                yield word + " "
                await asyncio.sleep(0.05)
            return

    # Redirection to the Multi-Provider Orchestration System (AI Router)
    from app.services.ai_router import ai_router

    keys = {
        "gemini_key": gemini_key,
        "openai_key": openai_key,
        "cohere_key": cohere_key,
        "anthropic_key": anthropic_key,
        "deepseek_key": deepseek_key
    }

    async for chunk in ai_router.stream_orchestrated_response(
        query=query,
        chat_history=chat_history,
        chat_mode=chat_mode,
        system_instructions=system_instructions,
        temperature=temperature,
        keys=keys,
        attachments=attachments
    ):
        yield chunk
    return

